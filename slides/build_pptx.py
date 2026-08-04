"""slides-spec.json からシニア向けセミナー用の PowerPoint を生成する。

    python build_pptx.py            # 生成
    python build_pptx.py --verify   # 生成物を再オープンして機械検証
    python build_pptx.py --all      # 生成 → 検証
"""

import argparse
import json
import math
import re
import sys
import unicodedata
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

BASE = Path(__file__).resolve().parent
SPEC_PATH = BASE / "slides-spec.json"
OUT_PATH = BASE / "deepfake-talk.pptx"

# ---------------------------------------------------------------- レイアウト定数
SLIDE_W, SLIDE_H = 13.333, 7.5
ML = 0.75                       # 左右マージン
CW = SLIDE_W - ML * 2           # 本文幅 11.833
INSET = 0.2                     # テキストボックス左右インセット合計

TITLE_TOP, TITLE_H = 0.45, 1.05  # 見出し（下揃え → 下端 1.50）
RULE_Y, RULE_W, RULE_H = 1.60, 2.6, 0.055
CONTENT_TOP = 1.80
CONTENT_BOTTOM = 6.55
FOOTER_Y, FOOTER_H = 6.40, 0.38
META_Y, META_H = 6.90, 0.40

BIG_TOP = 1.30
BIG_L, BIG_W = 0.55, 12.233

MUTED_ON_DARK = "9AA3B2"        # 濃色背景での補助色（#5C6472 は暗すぎる）
SUB_ON_DARK = "BCBDC0"          # 白 70% を #20242C 上で合成した値
EMPH_BG = "F3E6CF"
RULE_LIGHT = "E6E1D6"

META_NAMES = ("meta:pageno", "meta:part")
SLOT_NAMES = ("slot:image",)

# image_slot（あとから画像を貼るための予約領域）
SLOT_BG = "F1F2F5"              # ごく薄い背景
SLOT_LINE = "5C6472"            # 破線ボーダー（muted 系）
SLOT_LABEL = "イメージ画像スペース"
SLOT_LABEL_PT, SLOT_DESC_PT = 15, 12
SLOT_MAIN_R = 0.62              # pos:right のときの本文幅（CW 比）
SLOT_SIDE_R = 0.34              # 同・プレースホルダ幅
SLOT_GAP_R = 0.04               # 同・本文とのすき間
SLOT_BAND_H, SLOT_BAND_GAP = 1.35, 0.22   # pos:bottom の帯高さとすき間

# sources（出典）: 表示名＋URL の2行1組。比率は収まり判定と描画で共用する。
SRC_NAME_LS, SRC_URL_LS, SRC_GAP_R = 1.25, 1.15, 0.30
SRC_SIZES = (15, 14, 13, 12)    # 表示名の候補pt（15pt が全名を1行に収める上限）


# ---------------------------------------------------------------- 文字幅の見積り
def em_width(text):
    """全角=1.0em / 半角=0.55em として文字列の幅を em で返す。"""
    w = 0.0
    for ch in text:
        w += 1.0 if unicodedata.east_asian_width(ch) in ("W", "F", "A") else 0.55
    return w


def est_lines(text, size_pt, width_in):
    """size_pt の文字を width_in に流し込んだときの行数見積り。"""
    per_line = width_in * 72.0 / size_pt
    if per_line <= 0:
        return 1
    return max(1, math.ceil(em_width(text) / per_line))


# ---------------------------------------------------------------- 図形/文字ヘルパ
def flat(shape):
    """既定の枠線と影を落として、フラットな塗りだけの図形にする。"""
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_rect(slide, l, t, w, h, fill_hex, shape=MSO_SHAPE.RECTANGLE):
    shp = slide.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor.from_string(fill_hex)
    return flat(shp)


def add_tb(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP, margins=(0.1, 0.05)):
    """word_wrap 有効・自動リサイズ無効のテキストボックスを追加する。"""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(margins[0])
    tf.margin_top = tf.margin_bottom = Inches(margins[1])
    return box, tf


def _set_ea(run, font_name):
    """日本語が確実に指定フォントで出るよう a:ea / a:cs も設定する。"""
    rPr = run._r.get_or_add_rPr()
    for tag, successors in (
        ("a:ea", ("a:cs", "a:sym", "a:hlinkClick", "a:hlinkMouseOver", "a:rtl", "a:extLst")),
        ("a:cs", ("a:sym", "a:hlinkClick", "a:hlinkMouseOver", "a:rtl", "a:extLst")),
    ):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.insert_element_before(el, *successors)
        el.set("typeface", font_name)


def style_run(run, size, color, bold=False, serif=False):
    name = FONT_SERIF if serif else FONT_SANS
    f = run.font
    f.name = name
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = RGBColor.from_string(color)
    _set_ea(run, name)
    return run


def put_text(tf, text, size, color, bold=False, serif=False,
             align=PP_ALIGN.LEFT, line_spacing=None, first=True):
    """テキストフレームに1段落分（改行は段落分割）流し込む。"""
    for i, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if (first and i == 0) else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        style_run(p.add_run(), size, color, bold, serif).text = line
    return tf


def label_box(slide, l, t, w, h, text, size, color, bold=False, serif=False,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=None,
              margins=(0.1, 0.05), name=None):
    box, tf = add_tb(slide, l, t, w, h, anchor=anchor, margins=margins)
    if name:
        box.name = name
    put_text(tf, text, size, color, bold, serif, align, line_spacing)
    return box


def stack_rows(area_top, area_bottom, heights, gap_max=0.40):
    """行の高さリストから、上下中央寄せで各行の y 座標を返す。"""
    n = len(heights)
    total = sum(heights)
    room = area_bottom - area_top
    gap = min(gap_max, max(0.06, (room - total) / max(1, n - 1))) if n > 1 else 0.0
    block = total + gap * (n - 1)
    y = area_top + max(0.0, (room - block) / 2.0)
    ys = []
    for h in heights:
        ys.append(y)
        y += h + gap
    return ys


# ---------------------------------------------------------------- 共通パーツ
def add_meta(slide, n, part, dark=False):
    """右下のページ番号と左下の part。"""
    muted = MUTED_ON_DARK if dark else MUTED
    label_box(slide, SLIDE_W - ML - 1.2, META_Y, 1.2, META_H, str(n), 12, muted,
              align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, name="meta:pageno")
    if part:
        label_box(slide, ML, META_Y, 6.0, META_H, part, 12, muted,
                  anchor=MSO_ANCHOR.MIDDLE, name="meta:part")


def add_notes(slide, notes):
    slide.notes_slide.notes_text_frame.text = notes or ""


def add_title(slide, title):
    """アクセント色の見出し＋下線。見出しは 36pt 下限を守る。"""
    size = 40 if em_width(title) <= 18 else 36
    label_box(slide, ML, TITLE_TOP, CW, TITLE_H, title, size, ACCENT, bold=True,
              anchor=MSO_ANCHOR.BOTTOM)
    add_rect(slide, ML, RULE_Y, RULE_W, RULE_H, ACCENT)


def add_footer(slide, footer):
    label_box(slide, ML, FOOTER_Y, CW, FOOTER_H, footer, 15, MUTED)


def add_emph_box(slide, top, height, text, size=26, left=ML, width=CW):
    """アクセント背景の角丸ボックス（強調文）。"""
    shp = add_rect(slide, left, top, width, height, EMPH_BG, MSO_SHAPE.ROUNDED_RECTANGLE)
    try:
        shp.adjustments[0] = 0.16
    except (IndexError, ValueError):
        pass
    tf = shp.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.3)
    tf.margin_top = tf.margin_bottom = Inches(0.05)
    put_text(tf, text, size, INK, bold=True, align=PP_ALIGN.CENTER)
    return shp


def image_slot(s, pos):
    """指定位置の image_slot を返す（無ければ None）。"""
    slot = s.get("image_slot")
    return slot if slot and slot.get("pos") == pos else None


def add_image_slot(slide, l, t, w, h, desc):
    """あとから画像を貼るための予約領域を、破線の仮置き枠として描く。"""
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(l), Inches(t), Inches(w), Inches(h))
    shp.name = SLOT_NAMES[0]
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor.from_string(SLOT_BG)
    shp.line.color.rgb = RGBColor.from_string(SLOT_LINE)
    shp.line.width = Pt(0.75)
    shp.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    shp.shadow.inherit = False
    try:
        # 角丸の半径は min(w, h) 比。細長い枠でも 0.14in 程度に留める。
        shp.adjustments[0] = min(0.16, 0.14 / min(w, h))
    except (IndexError, ValueError):
        pass

    tf = shp.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.18)
    tf.margin_top = tf.margin_bottom = Inches(0.08)
    put_text(tf, SLOT_LABEL, SLOT_LABEL_PT, MUTED, bold=True,
             align=PP_ALIGN.CENTER, line_spacing=1.3)
    if desc:
        put_text(tf, desc, SLOT_DESC_PT, MUTED, align=PP_ALIGN.CENTER,
                 line_spacing=1.35, first=False)
    return shp


# ---------------------------------------------------------------- 各タイプの描画
def render_cover(slide, s):
    label_box(slide, ML, 2.30, CW, 1.80, s["title"], 48, INK, bold=True, serif=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if s.get("sub"):
        label_box(slide, ML, 4.25, CW, 0.85, s["sub"], 28, ACCENT,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if s.get("footer"):
        label_box(slide, ML, 6.10, CW, 0.55, s["footer"], 16, MUTED,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def render_section(slide, s):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, SECTION_BG)
    if s.get("kicker"):
        label_box(slide, ML, 0.65, 5.0, 0.60, s["kicker"], 20, ACCENT, bold=True,
                  anchor=MSO_ANCHOR.MIDDLE)
    title = s["title"]
    size = 54 if em_width(title) * 54 / 72 <= CW - INSET else 44
    label_box(slide, ML, 2.45, CW, 1.70, title, size, "FFFFFF", bold=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if s.get("sub"):
        label_box(slide, ML, 4.35, CW, 0.95, s["sub"], 24, SUB_ON_DARK,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def render_big(slide, s):
    bottom = CONTENT_BOTTOM
    if s.get("warn"):
        band_h = 1.05
        band_top = bottom - band_h
        bottom = band_top - 0.25
    slot = image_slot(s, "bottom")
    if slot:
        slot_top = bottom - SLOT_BAND_H
        bottom = slot_top - SLOT_BAND_GAP
    if s.get("sub"):
        sub_h = 1.10
        sub_top = bottom - sub_h
        bottom = sub_top - 0.15

    lines = str(s["big"]).split("\n")
    max_em = max(em_width(x) for x in lines)
    avail_w = BIG_W - INSET
    avail_h = bottom - BIG_TOP
    size = 44
    for cand in range(54, 43, -1):
        if (max_em * cand / 72 <= avail_w
                and len(lines) * cand * 1.28 / 72 <= avail_h):
            size = cand
            break
    box, tf = add_tb(slide, BIG_L, BIG_TOP, BIG_W, avail_h, anchor=MSO_ANCHOR.MIDDLE)
    put_text(tf, s["big"], size, INK, bold=True, serif=True,
             align=PP_ALIGN.CENTER, line_spacing=1.28)

    if s.get("sub"):
        label_box(slide, ML, sub_top, CW, sub_h, s["sub"], 22, MUTED,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.3)
    if slot:
        add_image_slot(slide, ML, slot_top, CW, SLOT_BAND_H, slot.get("desc"))
    if s.get("warn"):
        band = add_rect(slide, 0, band_top, SLIDE_W, band_h, DANGER)
        tf = band.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = Inches(0.5)
        put_text(tf, s["warn"], 28, "FFFFFF", bold=True, align=PP_ALIGN.CENTER)


def render_bullets(slide, s):
    add_title(slide, s["title"])
    bottom = CONTENT_BOTTOM
    if s.get("footer"):
        add_footer(slide, s["footer"])
        bottom = FOOTER_Y - 0.08
    # プレースホルダは本文領域（emphasis を含む）と同じ高さ帯に置く。
    slot = image_slot(s, "right")
    body_w = CW
    if slot:
        body_w = CW * SLOT_MAIN_R
        add_image_slot(slide, ML + CW * (SLOT_MAIN_R + SLOT_GAP_R), CONTENT_TOP,
                       CW * SLOT_SIDE_R, bottom - CONTENT_TOP, slot.get("desc"))
    emph_top = None
    if s.get("emphasis"):
        emph_h = 0.85
        if slot:
            # 幅が狭まって折り返す分、強調ボックスを縦に伸ばす（文字はみ出し防止）。
            emph_h = max(emph_h, est_lines(s["emphasis"], 26, body_w - 0.6)
                         * 26 * 1.35 / 72 + 0.22)
        emph_top = bottom - emph_h
        bottom = emph_top - 0.18

    items = s["bullets"]
    avail_h = bottom - CONTENT_TOP
    size = 24
    fit = False
    for cand in (tuple(range(28, 17, -1)) if slot else (28, 26, 24)):
        text_w = body_w - INSET - cand / 72.0  # ぶら下げインデント分を差し引く
        lines = sum(est_lines(t, cand, text_w) for t in items)
        need = lines * cand * 1.35 / 72 + (len(items) - 1) * (0.35 * cand / 72)
        if need <= avail_h:
            size, fit = cand, True
            break
    if slot and not fit:
        raise ValueError("n=%s: バレットが幅 %.2fin の本文領域に収まらない" % (s["n"], body_w))

    box, tf = add_tb(slide, ML, CONTENT_TOP, body_w, avail_h)
    hang = Pt(size)
    for i, text in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.35
        if i < len(items) - 1:
            p.space_after = Pt(size * 0.35)
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", str(hang.emu))
        pPr.set("indent", str(-hang.emu))
        style_run(p.add_run(), size, INK).text = "・" + text

    if emph_top is not None:
        add_emph_box(slide, emph_top, emph_h, s["emphasis"], 26, width=body_w)


_VALUE_ACCENT = r"(「[^」]*」|『[^』]*』|約?[0-9０-９][0-9０-９,，\.]*[万億兆]?円?)"


def _value_runs(p, value, size):
    """値の中の金額・鍵カッコ語をアクセント色、それ以外を本文色にする。"""
    for part in [x for x in re.split(_VALUE_ACCENT, value) if x]:
        hit = bool(re.fullmatch(_VALUE_ACCENT, part))
        style_run(p.add_run(), size, ACCENT if hit else INK, bold=True).text = part


def render_stat(slide, s):
    add_title(slide, s["title"])
    bottom = CONTENT_BOTTOM
    if s.get("footer"):
        add_footer(slide, s["footer"])
        bottom = FOOTER_Y - 0.08
    if s.get("emphasis"):
        emph_h = 0.70
        emph_top = bottom - emph_h
        bottom = emph_top - 0.15
        label_box(slide, ML, emph_top, CW, emph_h, s["emphasis"], 30, DANGER,
                  bold=True, anchor=MSO_ANCHOR.MIDDLE)

    rows = s["stats"]
    lab_w, val_x, val_w = 3.70, 4.75, 7.85
    heights = []
    for label, value in rows:
        h = max(est_lines(label, 20, lab_w - INSET) * 20 * 1.25 / 72,
                est_lines(value, 28, val_w - INSET) * 28 * 1.20 / 72) + 0.24
        heights.append(h)
    ys = stack_rows(CONTENT_TOP, bottom, heights, gap_max=0.20)

    for i, ((label, value), y, h) in enumerate(zip(rows, ys, heights)):
        label_box(slide, ML, y, lab_w, h, label, 20, MUTED,
                  anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.25)
        box, tf = add_tb(slide, val_x, y, val_w, h, anchor=MSO_ANCHOR.MIDDLE)
        tf.paragraphs[0].line_spacing = 1.2
        _value_runs(tf.paragraphs[0], value, 28)
        if i < len(rows) - 1:
            add_rect(slide, ML, y + h + 0.02, CW, 0.014, RULE_LIGHT)


def render_steps(slide, s):
    add_title(slide, s["title"])
    bottom = CONTENT_BOTTOM
    if s.get("footer"):
        add_footer(slide, s["footer"])
        bottom = FOOTER_Y - 0.08
    emph_top = None
    if s.get("emphasis"):
        emph_h = 0.85
        emph_top = bottom - emph_h
        bottom = emph_top - 0.18

    rows = s["steps"]
    wide_label = max(em_width(r[0]) for r in rows) > 5
    if wide_label:
        lab_w, main_x = 4.30, 5.35
    else:
        lab_w, main_x = 0.80, 1.65
    main_w = SLIDE_W - ML - main_x
    avail_h = bottom - CONTENT_TOP

    def row_heights(main_s, note_s):
        hs = []
        for r in rows:
            h = est_lines(r[1], main_s, main_w - INSET) * main_s * 1.20 / 72
            if len(r) > 2 and r[2]:
                h += est_lines(r[2], note_s, main_w - INSET) * note_s * 1.30 / 72 + 0.06
            hs.append(h)
        return hs

    combos = [(m, n) for m in (30, 28, 26) for n in (20, 19, 18)]
    main_s, note_s = combos[-1]
    heights = row_heights(main_s, note_s)
    for cand_main, cand_note in combos:
        hs = row_heights(cand_main, cand_note)
        if sum(hs) + (len(rows) - 1) * 0.14 <= avail_h:
            main_s, note_s, heights = cand_main, cand_note, hs
            break

    ys = stack_rows(CONTENT_TOP, bottom, heights, gap_max=0.35)
    lab_s = 20 if wide_label else 28
    lab_align = PP_ALIGN.LEFT if wide_label else PP_ALIGN.RIGHT

    for r, y, h in zip(rows, ys, heights):
        main_h = est_lines(r[1], main_s, main_w - INSET) * main_s * 1.20 / 72
        label_box(slide, ML if wide_label else ML + 0.05, y, lab_w, main_h, r[0],
                  lab_s, ACCENT, bold=True, align=lab_align,
                  anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
        label_box(slide, main_x, y, main_w, main_h, r[1], main_s, INK, bold=True,
                  anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)
        if len(r) > 2 and r[2]:
            label_box(slide, main_x, y + main_h + 0.04, main_w, h - main_h - 0.04,
                      r[2], note_s, MUTED, line_spacing=1.3)

    if emph_top is not None:
        add_emph_box(slide, emph_top, 0.85, s["emphasis"], 26)


def render_chart(slide, s):
    """見出し＋グラフ画像（アスペクト維持で本文領域に最大化・中央寄せ）＋footer。"""
    add_title(slide, s["title"])
    bottom = CONTENT_BOTTOM
    if s.get("footer"):
        add_footer(slide, s["footer"])
        bottom = FOOTER_Y - 0.08

    path = BASE / s["image"]
    if not path.exists():
        raise FileNotFoundError(
            "グラフ画像がない: %s（python build_chart.py で生成する）" % path)

    pic = slide.shapes.add_picture(str(path), Inches(0), Inches(0))
    avail_w, avail_h = Inches(CW), Inches(bottom - CONTENT_TOP)
    scale = min(avail_w / pic.width, avail_h / pic.height)
    pic.width, pic.height = int(pic.width * scale), int(pic.height * scale)
    pic.left = Emu(int(Inches(ML) + (avail_w - pic.width) / 2))
    pic.top = Emu(int(Inches(CONTENT_TOP) + (avail_h - pic.height) / 2))


def source_sizes(items, avail_h):
    """(表示名, URL) の2行1組を avail_h に収める (表示名pt, URLpt) を返す。"""
    size, url_size = SRC_SIZES[-1], 10
    for cand in SRC_SIZES:
        u = 11 if cand >= 14 else 10
        text_w = CW - INSET - cand / 72.0      # ぶら下げインデント分を差し引く
        name_lines = sum(est_lines(n, cand, text_w) for n, _ in items)
        url_lines = sum(est_lines(url, u, text_w) for _, url in items)
        need = (name_lines * cand * SRC_NAME_LS + url_lines * u * SRC_URL_LS
                + (len(items) - 1) * cand * SRC_GAP_R) / 72
        if need <= avail_h:
            size, url_size = cand, u
            break
    return size, url_size


def render_source_links(slide, items, avail_h):
    """出典を「表示名／URL」の2行1組で並べ、どちらの行にもリンクを張る。"""
    size, url_size = source_sizes(items, avail_h)
    box, tf = add_tb(slide, ML, CONTENT_TOP, CW, avail_h)
    hang = Pt(size)
    for i, (name, url) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = SRC_NAME_LS
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", str(hang.emu))
        pPr.set("indent", str(-hang.emu))
        run = style_run(p.add_run(), size, INK)
        run.text = "・" + name
        run.hyperlink.address = url

        q = tf.add_paragraph()
        q.line_spacing = SRC_URL_LS
        if i < len(items) - 1:
            q.space_after = Pt(size * SRC_GAP_R)
        qPr = q._p.get_or_add_pPr()
        qPr.set("marL", str(hang.emu))       # URL 行は「・」の分だけ字下げして表示名に揃える
        qPr.set("indent", "0")
        url_run = style_run(q.add_run(), url_size, MUTED)
        url_run.text = url
        url_run.hyperlink.address = url


def render_sources(slide, s):
    add_title(slide, s["title"])
    bottom = CONTENT_BOTTOM
    if s.get("footer"):
        add_footer(slide, s["footer"])
        bottom = FOOTER_Y - 0.08

    if s.get("sources"):
        render_source_links(slide, s["sources"], bottom - CONTENT_TOP)
        return

    items = s["bullets"]        # 旧形式（URL なしの文字列リスト）
    avail_h = bottom - CONTENT_TOP
    size = 16
    for cand in (18, 17, 16):
        text_w = CW - INSET - cand / 72.0
        lines = sum(est_lines(t, cand, text_w) for t in items)
        need = lines * cand * 1.30 / 72 + (len(items) - 1) * (0.45 * cand / 72)
        if need <= avail_h:
            size = cand
            break

    box, tf = add_tb(slide, ML, CONTENT_TOP, CW, avail_h)
    hang = Pt(size)
    for i, text in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.30
        if i < len(items) - 1:
            p.space_after = Pt(size * 0.45)
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", str(hang.emu))
        pPr.set("indent", str(-hang.emu))
        style_run(p.add_run(), size, INK).text = "・" + text


RENDERERS = {
    "cover": render_cover,
    "section": render_section,
    "big": render_big,
    "bullets": render_bullets,
    "stat": render_stat,
    "steps": render_steps,
    "chart": render_chart,
    "sources": render_sources,
}


# ---------------------------------------------------------------- ビルド
INK = MUTED = ACCENT = DANGER = SECTION_BG = "000000"
FONT_SANS = FONT_SERIF = "Yu Gothic"


def load_design(design):
    global INK, MUTED, ACCENT, DANGER, SECTION_BG, FONT_SANS, FONT_SERIF
    INK = design.get("ink", "20242C")
    MUTED = design.get("muted", "5C6472")
    ACCENT = design.get("accent", "B9791F")
    DANGER = design.get("danger", "A83B3B")
    SECTION_BG = design.get("section_bg", "20242C")
    FONT_SANS = design.get("font_sans", "Yu Gothic")
    FONT_SERIF = design.get("font_serif", "Yu Mincho")


def build():
    spec = json.loads(SPEC_PATH.read_text(encoding="utf-8"))
    load_design(spec["meta"].get("design", {}))

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]

    for s in spec["slides"]:
        slide = prs.slides.add_slide(blank)
        stype = s["type"]
        RENDERERS[stype](slide, s)
        add_meta(slide, s["n"], "" if stype == "cover" else s.get("part", ""),
                 dark=(stype == "section"))
        add_notes(slide, s.get("notes"))

    prs.save(str(OUT_PATH))
    return len(spec["slides"])


# ---------------------------------------------------------------- 検証
def verify():
    spec = json.loads(SPEC_PATH.read_text(encoding="utf-8"))
    expected = len(spec["slides"])
    prs = Presentation(str(OUT_PATH))
    ok = True

    n_slides = len(prs.slides)
    print("[1] slide count : %d (expected %d) -> %s"
          % (n_slides, expected, "OK" if n_slides == expected else "NG"))
    ok &= n_slides == expected

    bad_notes, bad_body, bad_font, fonts = [], [], [], set()
    for idx, slide in enumerate(prs.slides, start=1):
        # (b) 発表者ノート
        if not slide.has_notes_slide or not slide.notes_slide.notes_text_frame.text.strip():
            bad_notes.append(idx)
        # (c) meta 図形とプレースホルダを除いた本文テキスト
        body = 0
        for shp in slide.shapes:
            if not shp.has_text_frame or shp.name in META_NAMES + SLOT_NAMES:
                continue
            if shp.text_frame.text.strip():
                body += 1
        if body == 0:
            bad_body.append(idx)
        # (d) east asian フォント指定（meta 図形も含めて全 run を見る）
        for shp in slide.shapes:
            if not shp.has_text_frame:
                continue
            for p in shp.text_frame.paragraphs:
                for r in p.runs:
                    rPr = r._r.rPr
                    ea = None if rPr is None else rPr.find(qn("a:ea"))
                    if ea is None or ea.get("typeface") not in (FONT_SANS, FONT_SERIF):
                        bad_font.append(idx)
                    else:
                        fonts.add(ea.get("typeface"))

    print("[2] notes on all slides -> %s%s"
          % ("OK" if not bad_notes else "NG", "" if not bad_notes else " missing=%s" % bad_notes))
    print("[3] body text on all slides -> %s%s"
          % ("OK" if not bad_body else "NG", "" if not bad_body else " empty=%s" % bad_body))
    print("[4] a:ea typeface on all runs -> %s%s (typefaces=%s)"
          % ("OK" if not bad_font else "NG",
             "" if not bad_font else " bad=%s" % sorted(set(bad_font)),
             sorted(fonts)))
    ok &= not bad_notes and not bad_body and not bad_font

    # (e) chart スライスに画像が載っているか
    chart_idx = [i for i, s in enumerate(spec["slides"], start=1) if s["type"] == "chart"]
    no_pic = [i for i in chart_idx
              if not any(shp.shape_type == MSO_SHAPE_TYPE.PICTURE
                         for shp in prs.slides[i - 1].shapes)]
    e_ok = bool(chart_idx) and not no_pic
    print("[5] picture on chart slides : %s -> %s%s"
          % (chart_idx, "OK" if e_ok else "NG",
             "" if not no_pic else " missing=%s" % no_pic))
    ok &= e_ok

    # (f) 画像プレースホルダが image_slot 指定のスライドだけに載っているか
    want = [i for i, s in enumerate(spec["slides"], start=1) if s.get("image_slot")]
    got = [i for i, slide in enumerate(prs.slides, start=1)
           if any(shp.name in SLOT_NAMES for shp in slide.shapes)]
    f_ok = bool(want) and want == got
    print("[6] image placeholder slides : %s (expected %s) -> %s"
          % (got, want, "OK" if f_ok else "NG"))
    ok &= f_ok

    # (g) sources スライドの出典リンク（表示名行と URL 行の両方に spec の URL が張れているか）
    src_slides = [(i, s) for i, s in enumerate(spec["slides"], start=1) if s.get("sources")]
    bad_link, bad_order, n_runs, n_uniq = [], [], 0, 0
    for i, s in src_slides:
        want = {u for _, u in s["sources"]}
        got = set()
        for shp in prs.slides[i - 1].shapes:
            if not shp.has_text_frame:
                continue
            for p in shp.text_frame.paragraphs:
                for r in p.runs:
                    if not r.hyperlink.address:
                        continue
                    got.add(r.hyperlink.address)
                    n_runs += 1
                    # rPr の子要素順（a:ea/a:cs は a:hlinkClick より前）が崩れると
                    # PowerPoint が「修復が必要」と警告する。
                    tags = [c.tag for c in r._r.rPr]
                    h = tags.index(qn("a:hlinkClick"))
                    if any(qn(t) in tags and tags.index(qn(t)) > h for t in ("a:ea", "a:cs")):
                        bad_order.append(i)
        n_uniq += len(got)
        if got != want:
            bad_link.append(i)
    g_ok = not bad_link and not bad_order       # 旧 bullets 形式だけなら 0 件で OK
    print("[7] source hyperlinks on %s : %d unique (expected %d) / %d runs -> %s%s"
          % ([i for i, _ in src_slides], n_uniq,
             sum(len(s["sources"]) for _, s in src_slides), n_runs,
             "OK" if g_ok else "NG",
             "" if g_ok else " bad=%s order=%s" % (bad_link, sorted(set(bad_order)))))
    ok &= g_ok

    print("RESULT: %s" % ("PASS" if ok else "FAIL"))
    return 0 if ok else 1


def main():
    if hasattr(sys.stdout, "reconfigure"):
        sys.stdout.reconfigure(encoding="utf-8")
    ap = argparse.ArgumentParser()
    ap.add_argument("--verify", action="store_true", help="生成物を検証する")
    ap.add_argument("--all", action="store_true", help="生成してから検証する")
    args = ap.parse_args()

    if args.verify and not args.all:
        return verify()

    n = build()
    print("built: %s (%d slides)" % (OUT_PATH, n))
    if args.all:
        return verify()
    return 0


if __name__ == "__main__":
    sys.exit(main())
